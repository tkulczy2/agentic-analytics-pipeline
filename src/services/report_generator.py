"""PowerPoint report generation service."""
import logging
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from src.config import settings

logger = logging.getLogger(__name__)


class ReportGenerator:
    """Generate PowerPoint reports for MSSP analytics."""

    # Brand colors
    COLORS = {
        "primary": RGBColor(44, 62, 80),      # Dark blue
        "secondary": RGBColor(52, 152, 219),   # Light blue
        "success": RGBColor(40, 167, 69),      # Green
        "warning": RGBColor(255, 193, 7),      # Yellow
        "danger": RGBColor(220, 53, 69),       # Red
        "light": RGBColor(248, 249, 250),      # Light gray
        "dark": RGBColor(33, 37, 41),          # Dark gray
    }

    def __init__(self, reports_dir: Optional[str] = None):
        """Initialize the report generator."""
        self.reports_dir = Path(reports_dir or settings.reports_dir)
        self.reports_dir.mkdir(parents=True, exist_ok=True)

    def generate_executive_report(
        self,
        workflow_id: str,
        contract_id: str,
        performance_year: int,
        performance_month: int,
        financial_metrics: Dict[str, Any],
        quality_metrics: Dict[str, Any],
        risk_metrics: Dict[str, Any],
        predictions: Dict[str, Any],
    ) -> Path:
        """
        Generate executive summary PowerPoint report.

        Returns:
            Path to the generated report file
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title
        self._add_title_slide(
            prs, contract_id, performance_year, performance_month
        )

        # Slide 2: Executive Summary
        self._add_executive_summary_slide(
            prs, financial_metrics, quality_metrics, predictions
        )

        # Slide 3: Financial Performance
        self._add_financial_slide(prs, financial_metrics)

        # Slide 4: Quality Performance
        self._add_quality_slide(prs, quality_metrics)

        # Slide 5: Risk Stratification
        self._add_risk_slide(prs, risk_metrics)

        # Slide 6: Key Insights
        self._add_insights_slide(prs, predictions)

        # Slide 7: Predictions & Recommendations
        self._add_predictions_slide(prs, predictions, financial_metrics)

        # Save the presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{contract_id}_{performance_year}_M{performance_month:02d}_{timestamp}.pptx"
        filepath = self.reports_dir / filename

        prs.save(str(filepath))
        logger.info(f"Generated executive report: {filepath}")

        return filepath

    def _add_title_slide(
        self,
        prs: Presentation,
        contract_id: str,
        performance_year: int,
        performance_month: int
    ):
        """Add title slide."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS["primary"]
        background.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "MSSP Analytics Report"
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = f"{contract_id} | {performance_year} M{performance_month:02d}"
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = self.COLORS["secondary"]
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Date
        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        date_para.font.size = Pt(16)
        date_para.font.color.rgb = RGBColor(200, 200, 200)
        date_para.alignment = PP_ALIGN.CENTER

    def _add_executive_summary_slide(
        self,
        prs: Presentation,
        financial: Dict[str, Any],
        quality: Dict[str, Any],
        predictions: Dict[str, Any]
    ):
        """Add executive summary slide with key metrics."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title
        self._add_slide_title(slide, "Executive Summary")

        # Metric boxes
        metrics = [
            {
                "label": "Total Savings",
                "value": f"${financial.get('total_savings', 0):,.0f}",
                "subtext": f"{financial.get('savings_percentage', 0):.1f}% below baseline",
                "color": self.COLORS["success"] if financial.get('total_savings', 0) > 0 else self.COLORS["danger"],
            },
            {
                "label": "Quality Score",
                "value": f"{quality.get('composite_score', 0):.1f}%",
                "subtext": f"Gate: {quality.get('quality_gate_status', 'pending').title()}",
                "color": self.COLORS["success"] if quality.get('quality_gate_status') == 'eligible' else self.COLORS["warning"],
            },
            {
                "label": "Projected Payment",
                "value": f"${predictions.get('projected_shared_savings', 0):,.0f}",
                "subtext": f"{predictions.get('probability_shared_savings', 0)*100:.0f}% probability",
                "color": self.COLORS["secondary"],
            },
            {
                "label": "Members",
                "value": f"{financial.get('average_members', 0):,}",
                "subtext": f"{financial.get('member_months', 0):,} member-months",
                "color": self.COLORS["primary"],
            },
        ]

        x_start = Inches(0.5)
        y_start = Inches(1.8)
        box_width = Inches(3)
        box_height = Inches(2)
        gap = Inches(0.2)

        for i, metric in enumerate(metrics):
            x = x_start + i * (box_width + gap)
            self._add_metric_box(
                slide, x, y_start, box_width, box_height,
                metric["label"], metric["value"], metric["subtext"], metric["color"]
            )

        # Status summary
        status_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5), Inches(12.333), Inches(1.5)
        )
        status_frame = status_box.text_frame
        status_frame.word_wrap = True

        risks = predictions.get("risks", [])
        opportunities = predictions.get("opportunities", [])

        p = status_frame.paragraphs[0]
        p.text = f"Key Findings: {len(risks)} risks identified, {len(opportunities)} opportunities"
        p.font.size = Pt(18)
        p.font.color.rgb = self.COLORS["dark"]

    def _add_financial_slide(self, prs: Presentation, financial: Dict[str, Any]):
        """Add financial performance slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Financial Performance")

        # Create spending comparison data
        baseline = financial.get("baseline_spending", 0)
        actual = financial.get("actual_spending", 0)
        target = baseline * (1 - financial.get("target_reduction_pct", 0.05))

        # Simple bar representation using shapes
        bar_width = Inches(2)
        max_height = Inches(4)
        y_bottom = Inches(6)
        x_start = Inches(2)
        gap = Inches(1.5)

        max_value = max(baseline, actual, target) or 1

        bars = [
            ("Baseline", baseline, self.COLORS["primary"]),
            ("Target", target, self.COLORS["secondary"]),
            ("Actual", actual, self.COLORS["success"] if actual < baseline else self.COLORS["danger"]),
        ]

        for i, (label, value, color) in enumerate(bars):
            x = x_start + i * (bar_width + gap)
            height = (value / max_value) * max_height
            y = y_bottom - height

            # Bar
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, bar_width, height
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()

            # Label
            label_box = slide.shapes.add_textbox(x, y_bottom + Inches(0.1), bar_width, Inches(0.4))
            label_frame = label_box.text_frame
            p = label_frame.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # Value
            value_box = slide.shapes.add_textbox(x, y - Inches(0.4), bar_width, Inches(0.4))
            value_frame = value_box.text_frame
            p = value_frame.paragraphs[0]
            p.text = f"${value/1000000:.1f}M"
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER

        # PMPM metrics
        pmpm_box = slide.shapes.add_textbox(Inches(9), Inches(2), Inches(4), Inches(4))
        pmpm_frame = pmpm_box.text_frame
        pmpm_frame.word_wrap = True

        p = pmpm_frame.paragraphs[0]
        p.text = "PMPM Metrics"
        p.font.bold = True
        p.font.size = Pt(18)

        metrics_text = [
            f"Baseline PMPM: ${financial.get('baseline_pmpm', 0):,.0f}",
            f"Target PMPM: ${financial.get('target_pmpm', 0):,.0f}",
            f"Actual PMPM: ${financial.get('actual_pmpm', 0):,.0f}",
            "",
            f"ER Visits/1000: {financial.get('er_visits_per_1000', 0):.0f}",
            f"Admits/1000: {financial.get('admits_per_1000', 0):.0f}",
        ]

        for text in metrics_text:
            p = pmpm_frame.add_paragraph()
            p.text = text
            p.font.size = Pt(14)

    def _add_quality_slide(self, prs: Presentation, quality: Dict[str, Any]):
        """Add quality performance slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Quality Performance")

        # Composite score
        score = quality.get("composite_score", 0)
        threshold = quality.get("quality_threshold", 80)
        status = quality.get("quality_gate_status", "pending")

        score_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(4), Inches(1.5)
        )
        score_frame = score_box.text_frame

        p = score_frame.paragraphs[0]
        p.text = f"{score:.1f}%"
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = (
            self.COLORS["success"] if status == "eligible"
            else self.COLORS["warning"] if status == "at_risk"
            else self.COLORS["danger"]
        )

        p = score_frame.add_paragraph()
        p.text = f"Composite Score (Threshold: {threshold}%)"
        p.font.size = Pt(16)

        p = score_frame.add_paragraph()
        p.text = f"Quality Gate: {status.upper()}"
        p.font.size = Pt(18)
        p.font.bold = True

        # Category scores
        categories = [
            ("Preventive Care", quality.get("preventive_care_score", 0), "1.0x"),
            ("Chronic Disease", quality.get("chronic_disease_score", 0), "2.0x"),
            ("Care Coordination", quality.get("care_coordination_score", 0), "1.5x"),
            ("Patient Experience", quality.get("patient_experience_score", 0), "1.0x"),
        ]

        y_start = Inches(2)
        bar_max_width = Inches(6)
        bar_height = Inches(0.5)
        x_start = Inches(5.5)

        for i, (name, score_val, weight) in enumerate(categories):
            y = y_start + i * Inches(1.2)

            # Label
            label_box = slide.shapes.add_textbox(x_start, y - Inches(0.3), Inches(3), Inches(0.3))
            label_frame = label_box.text_frame
            p = label_frame.paragraphs[0]
            p.text = f"{name} ({weight})"
            p.font.size = Pt(12)

            # Background bar
            bg_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x_start, y, bar_max_width, bar_height
            )
            bg_bar.fill.solid()
            bg_bar.fill.fore_color.rgb = self.COLORS["light"]
            bg_bar.line.fill.background()

            # Progress bar
            progress_width = (score_val / 100) * bar_max_width
            if progress_width > 0:
                progress_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, x_start, y, progress_width, bar_height
                )
                progress_bar.fill.solid()
                progress_bar.fill.fore_color.rgb = (
                    self.COLORS["success"] if score_val >= 80
                    else self.COLORS["warning"] if score_val >= 70
                    else self.COLORS["danger"]
                )
                progress_bar.line.fill.background()

            # Score value
            value_box = slide.shapes.add_textbox(
                x_start + bar_max_width + Inches(0.2), y, Inches(1), bar_height
            )
            value_frame = value_box.text_frame
            p = value_frame.paragraphs[0]
            p.text = f"{score_val:.1f}%"
            p.font.size = Pt(14)
            p.font.bold = True

    def _add_risk_slide(self, prs: Presentation, risk: Dict[str, Any]):
        """Add risk stratification slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Risk Stratification")

        # Risk distribution
        total = risk.get("total_members", 1)
        categories = [
            ("Low Risk", risk.get("low_risk_count", 0), risk.get("low_risk_pmpm", 0), self.COLORS["success"]),
            ("Medium Risk", risk.get("medium_risk_count", 0), risk.get("medium_risk_pmpm", 0), self.COLORS["warning"]),
            ("High Risk", risk.get("high_risk_count", 0), risk.get("high_risk_pmpm", 0), self.COLORS["danger"]),
        ]

        x_start = Inches(0.5)
        y_start = Inches(2)
        box_width = Inches(4)
        box_height = Inches(1.5)

        for i, (name, count, pmpm, color) in enumerate(categories):
            y = y_start + i * (box_height + Inches(0.3))
            pct = (count / total * 100) if total > 0 else 0

            # Background
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x_start, y, box_width, box_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()

            # Text
            text_box = slide.shapes.add_textbox(x_start + Inches(0.2), y + Inches(0.2), box_width - Inches(0.4), box_height - Inches(0.4))
            text_frame = text_box.text_frame

            p = text_frame.paragraphs[0]
            p.text = name
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

            p = text_frame.add_paragraph()
            p.text = f"{count:,} members ({pct:.1f}%)"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)

            p = text_frame.add_paragraph()
            p.text = f"PMPM: ${pmpm:,.0f}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)

        # Average risk score
        avg_box = slide.shapes.add_textbox(Inches(6), Inches(2), Inches(6), Inches(3))
        avg_frame = avg_box.text_frame

        p = avg_frame.paragraphs[0]
        p.text = "Average Risk Score"
        p.font.size = Pt(18)
        p.font.bold = True

        p = avg_frame.add_paragraph()
        p.text = f"{risk.get('average_risk_score', 1.0):.2f}"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["primary"]

        p = avg_frame.add_paragraph()
        p.text = "(1.0 = Medicare average)"
        p.font.size = Pt(14)

    def _add_insights_slide(self, prs: Presentation, predictions: Dict[str, Any]):
        """Add key insights slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Key Insights")

        risks = predictions.get("risks", [])
        opportunities = predictions.get("opportunities", [])

        # Combine and take top 4
        all_items = [("risk", r) for r in risks] + [("opportunity", o) for o in opportunities]
        all_items = all_items[:4]

        if not all_items:
            all_items = [("info", {"title": "No significant findings", "description": "Performance is within expected ranges."})]

        x_positions = [Inches(0.5), Inches(6.7)]
        y_positions = [Inches(1.8), Inches(4.3)]
        box_width = Inches(6)
        box_height = Inches(2)

        for i, (item_type, item) in enumerate(all_items):
            x = x_positions[i % 2]
            y = y_positions[i // 2]

            color = (
                self.COLORS["danger"] if item_type == "risk"
                else self.COLORS["success"] if item_type == "opportunity"
                else self.COLORS["secondary"]
            )

            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.COLORS["light"]
            box.line.color.rgb = color

            # Icon indicator
            indicator = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.3), Inches(0.3)
            )
            indicator.fill.solid()
            indicator.fill.fore_color.rgb = color
            indicator.line.fill.background()

            # Text
            text_box = slide.shapes.add_textbox(
                x + Inches(0.6), y + Inches(0.15),
                box_width - Inches(0.8), box_height - Inches(0.3)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            p = text_frame.paragraphs[0]
            p.text = item.get("title", "Finding")
            p.font.size = Pt(16)
            p.font.bold = True

            p = text_frame.add_paragraph()
            p.text = item.get("description", "")
            p.font.size = Pt(12)

            if item.get("recommendation"):
                p = text_frame.add_paragraph()
                p.text = f"Recommendation: {item.get('recommendation')}"
                p.font.size = Pt(11)
                p.font.italic = True

    def _add_predictions_slide(
        self,
        prs: Presentation,
        predictions: Dict[str, Any],
        financial: Dict[str, Any]
    ):
        """Add predictions and recommendations slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Year-End Projections")

        # Projections
        proj_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(4))
        proj_frame = proj_box.text_frame
        proj_frame.word_wrap = True

        p = proj_frame.paragraphs[0]
        p.text = "Projected Year-End Results"
        p.font.size = Pt(20)
        p.font.bold = True

        projected_savings = predictions.get("projected_year_end_savings", 0)
        projected_shared = predictions.get("projected_shared_savings", 0)
        prob_savings = predictions.get("probability_shared_savings", 0)

        items = [
            f"Projected Savings: ${projected_savings:,.0f}",
            f"Projected Shared Savings: ${projected_shared:,.0f}",
            f"Probability of Achieving Target: {prob_savings*100:.0f}%",
            "",
            f"Savings Range (95% CI):",
            f"  Lower: ${predictions.get('savings_lower_bound', 0):,.0f}",
            f"  Upper: ${predictions.get('savings_upper_bound', 0):,.0f}",
        ]

        for item in items:
            p = proj_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)

        # Recommendations
        rec_box = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.8), Inches(4))
        rec_frame = rec_box.text_frame
        rec_frame.word_wrap = True

        p = rec_frame.paragraphs[0]
        p.text = "Recommended Actions"
        p.font.size = Pt(20)
        p.font.bold = True

        risks = predictions.get("risks", [])
        recommendations = [r.get("recommendation") for r in risks if r.get("recommendation")][:4]

        if not recommendations:
            recommendations = ["Continue current care management initiatives", "Monitor high-risk member engagement"]

        for i, rec in enumerate(recommendations, 1):
            p = rec_frame.add_paragraph()
            p.text = f"{i}. {rec}"
            p.font.size = Pt(14)

    def _add_slide_title(self, slide, title: str):
        """Add title to a slide."""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["primary"]

    def _add_metric_box(
        self,
        slide,
        x, y, width, height,
        label: str,
        value: str,
        subtext: str,
        color: RGBColor
    ):
        """Add a metric box to a slide."""
        # Background
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = self.COLORS["light"]
        box.line.color.rgb = color

        # Color bar at top
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.15)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Text
        text_box = slide.shapes.add_textbox(
            x + Inches(0.1), y + Inches(0.3),
            width - Inches(0.2), height - Inches(0.4)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = self.COLORS["dark"]
        p.alignment = PP_ALIGN.CENTER

        p = text_frame.add_paragraph()
        p.text = value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        p = text_frame.add_paragraph()
        p.text = subtext
        p.font.size = Pt(10)
        p.font.color.rgb = self.COLORS["dark"]
        p.alignment = PP_ALIGN.CENTER
